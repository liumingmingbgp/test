from pptx import Presentation
from pptx.util import Inches

ppt = Presentation()
layout = ppt.slide_layouts[6]
slide = ppt.slides.add_slide(layout)

### 定义图片位置大小 ###
left = Inches(0)
top = Inches(0)
width = Inches(3)
height = Inches(2)
img_path = 'ICON.jpeg'

pic = slide.shapes.add_picture(img_path, left, top, width, height)

ppt.save('6_PPT\Add_img.pptx')