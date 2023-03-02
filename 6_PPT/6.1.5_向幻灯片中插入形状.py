from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

ppt = Presentation()
layout = ppt.slide_layouts[6]
slide = ppt.slides.add_slide(layout)

### 定义插入形状的位置大小 ###
left = Inches(1)
top = Inches(2)
width = Inches(1.8)
height = Inches(1)

### 插入形状 ###
shape = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
shape.text = '第一步'

for i in range(2, 6):
    left = left + width - Inches(0.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.text = f'第{i}步'
ppt.save('6_PPT\Add_shape.pptx')