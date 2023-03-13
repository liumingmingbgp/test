from pptx import Presentation

ppt = Presentation()

for layout in ppt.slide_layouts:
    slide=ppt.slides.add_slide(layout)
ppt.save('6_PPT\show_all_layouts.pptx')