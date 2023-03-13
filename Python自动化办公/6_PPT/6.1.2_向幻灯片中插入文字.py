from pptx import Presentation

ppt = Presentation()

### 添加第一页幻灯片 ###
slide_layout = ppt.slide_layouts[0]
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = '我是标题'
subtitle = slide.placeholders[1] # 获取本页幻灯片的第二个文本框
subtitle.text = '正文框'

### 添加第二页幻灯片 ###
slide_layout = ppt.slide_layouts[1]
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = '我是标题2'
subtitle = slide.placeholders[1]
subtitle.text = '正文框2'
ppt.save('6_PPT\write_text.pptx')

### 通过占位符向ppt中插入文字 ###
layout = ppt.slide_layouts[0]
slide = ppt.slides.add_slide(layout)
placeholders = slide.shapes.placeholders
placeholders[0].text = '第一个文本框'
placeholders[1].text = '第二个文本框'
ppt.save('6_PPT\write_text2.pptx')

### 将新的文字追加到已有文字之后 ###
ppt = Presentation('6_PPT\write_text.pptx')
slide0 = ppt.slides[0]  # 获取第一页幻灯片
placeholder = slide0.shapes.placeholders # 获取第一页幻灯片中所有的占位符
new_paragraph = placeholder[1].text_frame.add_paragraph() # 在第二个占位符后加入内容
new_paragraph.text = '追加新的文字'
ppt.save('6_PPT\write_text3.pptx')