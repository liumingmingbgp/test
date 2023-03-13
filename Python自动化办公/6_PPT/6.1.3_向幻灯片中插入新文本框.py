from pptx import Presentation
from pptx.util import Inches

ppt = Presentation()
layout = ppt.slide_layouts[6]  # 添加空白slide
slide = ppt.slides.add_slide(layout)

### 预设文本框位置及大小 ###
left = Inches(5)
top = Inches(5)
width = Inches(5)
height = Inches(5)

textbox = slide.shapes.add_textbox(left, top, width, height)
textbox.text = '这是一个新的文本框'

### 添加新段落
new_paragraph = textbox.text_frame.add_paragraph()
new_paragraph.text = '第二段内容'

ppt.save('6_PPT\Add_new_text.pptx')