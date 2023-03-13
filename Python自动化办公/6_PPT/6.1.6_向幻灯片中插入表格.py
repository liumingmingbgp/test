from pptx import Presentation
from pptx.util import Inches

ppt = Presentation()
layout = ppt.slide_layouts[6]
slide = ppt.slides.add_slide(layout)

rows = 2
cols = 2
left = Inches(3.5)
top = Inches(4.5)
width = Inches(6)
height = Inches(0.8)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

table.cell(0,0).text = '第一行第一列'
table.cell(0,1).text = '第一行第二列'
table.cell(1,0).text = '第二行第一列'
table.cell(1,1).text = '第二行第二列'


ppt.save('6_PPT\Add_table.pptx')